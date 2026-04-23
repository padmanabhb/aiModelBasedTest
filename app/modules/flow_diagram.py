"""
Flow Diagram Generator Module
Creates UML and flow diagrams that can be exported to Visio and PowerPoint
"""
import json
from typing import Dict, List, Tuple, Any
from dataclasses import dataclass, asdict

@dataclass
class Node:
    """Represents a node in the flow diagram"""
    id: str
    label: str
    node_type: str = "process"  # process, decision, endpoint, startpoint
    x: float = 0
    y: float = 0

@dataclass
class Edge:
    """Represents an edge between nodes"""
    source: str
    target: str
    label: str = ""
    condition: str = ""

class FlowDiagramGenerator:
    """
    Generates optimized flow diagrams from requirement analysis
    """
    
    def __init__(self):
        """Initialize the diagram generator"""
        self.nodes: Dict[str, Node] = {}
        self.edges: List[Edge] = []
        self.paths: Dict[str, List[str]] = {}
    
    def create_diagram_from_flows(self, test_flows: List[Dict]) -> Dict[str, Any]:
        """
        Create a flow diagram from test flows with user navigation
        
        Args:
            test_flows: List of test flow dictionaries
            
        Returns:
            Dictionary containing nodes, edges, and diagram data
        """
        self.nodes = {}
        self.edges = []
        
        # First, create nodes from user_screens if available
        user_screens = test_flows[0].get('__user_screens__', []) if test_flows and isinstance(test_flows[0], dict) else []
        navigation_paths = test_flows[0].get('__navigation_paths__', []) if test_flows and isinstance(test_flows[0], dict) else []
        
        if user_screens:
            # Screen-based diagram
            return self._create_screen_based_diagram(user_screens, navigation_paths, test_flows)
        else:
            # Legacy flow-based diagram
            return self._create_flow_based_diagram(test_flows)
    
    def _create_screen_based_diagram(self, user_screens: List[Dict], navigation_paths: List[Dict], test_flows: List[Dict]) -> Dict[str, Any]:
        """
        Create a diagram based on user screens and navigation paths
        
        Args:
            user_screens: List of screen definitions
            navigation_paths: List of navigation paths between screens
            test_flows: List of test flows that use these screens
            
        Returns:
            Dictionary containing diagram data
        """
        self.nodes = {}
        self.edges = []
        
        # Color map for different screen types
        screen_colors = {
            "landing": "#28a745",      # Green
            "product": "#0066cc",      # Blue
            "list": "#007bff",         # Light Blue
            "detail": "#0099ff",       # Lighter Blue
            "form": "#ffc107",         # Amber
            "confirmation": "#28a745", # Green
        }
        
        # First pass: identify decision points (screens with multiple outgoing edges)
        outgoing_edges = {}
        for nav_path in navigation_paths:
            from_screen = nav_path.get('from_screen')
            if from_screen not in outgoing_edges:
                outgoing_edges[from_screen] = []
            outgoing_edges[from_screen].append(nav_path.get('to_screen'))
        
        decision_screens = {screen_id for screen_id, targets in outgoing_edges.items() if len(set(targets)) > 1}
        
        # Create nodes for each screen
        layout_rows = {}
        for i, screen in enumerate(user_screens):
            screen_id = screen.get('screen_id', f'SCREEN_{i}')
            screen_name = screen.get('screen_name', screen_id)
            screen_type = screen.get('screen_type', 'process')
            
            # Mark as decision if it's a divergence point
            if screen_id in decision_screens:
                node_type = "decision"
            elif screen_type == "landing":
                node_type = "startpoint"
            elif screen_type == "confirmation":
                node_type = "endpoint"
            else:
                node_type = "process"
            
            node = Node(
                screen_id,
                screen_name,
                node_type,
                x=100 + (i % 4) * 200,
                y=100 + (i // 4) * 150
            )
            self.nodes[screen_id] = node
            
            # Track for layout
            row = i // 4
            if row not in layout_rows:
                layout_rows[row] = []
            layout_rows[row].append(screen_id)
        
        # Create edges from navigation paths
        edge_set = set()
        for nav_path in navigation_paths:
            from_screen = nav_path.get('from_screen')
            to_screen = nav_path.get('to_screen')
            action = nav_path.get('action', '')
            flows = nav_path.get('flows', [])
            
            # For decision nodes, use the target screen name as the path label
            if from_screen in decision_screens:
                target_name = self.nodes[to_screen].label if to_screen in self.nodes else to_screen
                # Clean up the path label - use concise descriptive text
                path_label = target_name
            else:
                path_label = action if action else f"To {self.nodes[to_screen].label}" if to_screen in self.nodes else action
            
            # Create edge label showing flows that use this path (optional)
            flow_labels = ','.join([f.replace('UF', 'F') for f in flows[:2]])  # Limit label length
            edge_label = path_label
            
            edge_key = (from_screen, to_screen)
            if edge_key not in edge_set:
                self.edges.append(Edge(from_screen, to_screen, edge_label))
                edge_set.add(edge_key)
        
        # Return complete diagram data
        return self._generate_diagram_data()
    
    def _create_flow_based_diagram(self, test_flows: List[Dict]) -> Dict[str, Any]:
        """
        Create a legacy flow-based diagram
        
        Args:
            test_flows: List of test flow dictionaries
            
        Returns:
            Dictionary containing diagram data
        """
        self.nodes = {}
        self.edges = []
        
        # Create start node
        start_node = Node("START", "Start", "startpoint", x=50, y=50)
        self.nodes["START"] = start_node
        
        y_position = 100
        flow_nodes = {}
        
        # Create nodes for each flow
        for i, flow in enumerate(test_flows):
            flow_id = flow.get("flow_id", f"TF{i:02d}")
            node_id = f"FLOW_{flow_id}"
            node = Node(
                node_id,
                flow.get("name", f"Flow {i+1}"),
                "process",
                x=50 + (i % 3) * 200,
                y=y_position + (i // 3) * 100
            )
            self.nodes[node_id] = node
            flow_nodes[flow_id] = node_id
            
            # Create nodes for steps
            steps = flow.get("steps", [])
            for j, step in enumerate(steps):
                step_id = f"{flow_id}_STEP_{j}"
                # Use full step description without truncation
                step_description = step if len(step) <= 50 else step[:47] + "..."
                step_node = Node(
                    step_id,
                    step_description,
                    "process",
                    x=50 + (j % 2) * 150,
                    y=y_position + 150 + (j // 2) * 60
                )
                self.nodes[step_id] = step_node
        
        # Create end node
        end_node = Node("END", "End", "endpoint", x=50, y=y_position + 400)
        self.nodes["END"] = end_node
        
        # Connect flows
        for i, (flow_id, node_id) in enumerate(flow_nodes.items()):
            # Connect start to first flows
            if i == 0:
                self.edges.append(Edge("START", node_id, "Initialize"))
            
            # Connect between flows
            if i > 0:
                prev_node_id = list(flow_nodes.values())[i - 1]
                self.edges.append(Edge(prev_node_id, node_id, f"Flow {i}"))
            
            # Connect to end
            if i == len(flow_nodes) - 1:
                self.edges.append(Edge(node_id, "END", "Complete"))
        
        return self._generate_diagram_data()
    
    def create_diagram_from_entities(self, entities: List[str], 
                                     relationships: Dict[str, List[str]]) -> Dict[str, Any]:
        """
        Create a diagram from entities and relationships
        
        Args:
            entities: List of entity names
            relationships: Dictionary of entity relationships
            
        Returns:
            Dictionary containing the diagram
        """
        self.nodes = {}
        self.edges = []
        
        # Create nodes for entities
        positions = self._calculate_positions(len(entities))
        for i, entity in enumerate(entities):
            node = Node(
                f"ENTITY_{i}",
                entity,
                "process",
                x=positions[i][0],
                y=positions[i][1]
            )
            self.nodes[f"ENTITY_{i}"] = node
        
        # Create edges for relationships
        entity_to_id = {entity: f"ENTITY_{i}" for i, entity in enumerate(entities)}
        
        for source_entity, targets in relationships.items():
            if source_entity in entity_to_id:
                source_id = entity_to_id[source_entity]
                for target_entity in targets:
                    if target_entity in entity_to_id:
                        target_id = entity_to_id[target_entity]
                        self.edges.append(Edge(source_id, target_id, "relates to"))
        
        return self._generate_diagram_data()
    
    def _calculate_positions(self, count: int) -> List[Tuple[float, float]]:
        """Calculate node positions in a circular layout"""
        import math
        positions = []
        radius = 200
        center_x, center_y = 300, 300
        
        for i in range(count):
            angle = (2 * math.pi * i) / count
            x = center_x + radius * math.cos(angle)
            y = center_y + radius * math.sin(angle)
            positions.append((x, y))
        
        return positions
    
    def _generate_diagram_data(self) -> Dict[str, Any]:
        """Generate diagram data in exportable format"""
        nodes_data = []
        for node_id, node in self.nodes.items():
            nodes_data.append({
                "id": node.id,
                "label": node.label,
                "type": node.node_type,
                "position": {"x": node.x, "y": node.y}
            })
        
        edges_data = []
        for edge in self.edges:
            edges_data.append({
                "source": edge.source,
                "target": edge.target,
                "label": edge.label,
                "condition": edge.condition
            })
        
        return {
            "nodes": nodes_data,
            "edges": edges_data,
            "stats": {
                "total_nodes": len(self.nodes),
                "total_edges": len(self.edges),
                "node_types": self._count_node_types()
            }
        }
    
    def _count_node_types(self) -> Dict[str, int]:
        """Count nodes by type"""
        type_count = {}
        for node in self.nodes.values():
            type_count[node.node_type] = type_count.get(node.node_type, 0) + 1
        return type_count
    
    def export_to_mxgraph(self, filename: str) -> str:
        """
        Export diagram to mxGraph XML format (editable in draw.io, importable to Visio)
        
        Args:
            filename: Output filename
            
        Returns:
            XML content
        """
        import xml.sax.saxutils as saxutils
        
        xml = '<?xml version="1.0" encoding="UTF-8"?>\n'
        xml += '<mxfile host="app.diagrams.net" modified="2024-01-01T00:00:00" version="21">\n'
        xml += '  <diagram name="Page-1" id="page-1">\n'
        xml += '    <mxGraphModel dx="1200" dy="800" grid="1" gridSize="10" guides="1" tooltips="1" connect="1" arrows="1" fold="1" page="1" pageScale="1" pageWidth="1200" pageHeight="800" background="#ffffff" math="0" shadow="0">\n'
        xml += '      <root>\n'
        xml += '        <mxCell id="0" />\n'
        xml += '        <mxCell id="1" parent="0" />\n'
        
        # Add nodes
        for node in self.nodes.values():
            shape = self._get_shape_for_type(node.node_type)
            safe_label = saxutils.escape(node.label)
            xml += f'        <mxCell id="{node.id}" value="{safe_label}" style="{shape}" vertex="1" parent="1">\n'
            xml += f'          <mxGeometry x="{node.x}" y="{node.y}" width="100" height="60" as="geometry" />\n'
            xml += '        </mxCell>\n'
        
        # Add edges
        for i, edge in enumerate(self.edges):
            safe_label = saxutils.escape(edge.label)
            xml += f'        <mxCell id="edge_{i}" value="{safe_label}" style="edgeStyle=orthogonalEdgeStyle" edge="1" parent="1" source="{edge.source}" target="{edge.target}">\n'
            xml += '          <mxGeometry relative="1" as="geometry" />\n'
            xml += '        </mxCell>\n'
        
        xml += '      </root>\n'
        xml += '    </mxGraphModel>\n'
        xml += '  </diagram>\n'
        xml += '</mxfile>\n'
        
        return xml
    
    def _get_shape_for_type(self, node_type: str) -> str:
        """Get mxGraph style for node type"""
        styles = {
            "startpoint": "ellipse;whiteSpace=wrap;html=1;fillColor=#00FF00;",
            "endpoint": "ellipse;whiteSpace=wrap;html=1;fillColor=#FF0000;",
            "process": "rectangle;whiteSpace=wrap;html=1;fillColor=#0099FF;",
            "decision": "rhombus;whiteSpace=wrap;html=1;fillColor=#FFCC00;"
        }
        return styles.get(node_type, styles["process"])
    
    def export_to_json(self) -> str:
        """Export diagram as JSON"""
        data = self._generate_diagram_data()
        return json.dumps(data, indent=2)
    
    def generate_pptx_compatible_data(self) -> Dict[str, Any]:
        """
        Generate data compatible with python-pptx for PowerPoint export
        
        Returns:
            Dictionary with PowerPoint-compatible structure
        """
        slides = []
        
        # Slide 1: Overview
        slides.append({
            "title": "Flow Diagram Overview",
            "type": "title",
            "content": {
                "subtitle": f"Total Nodes: {len(self.nodes)}, Total Edges: {len(self.edges)}"
            }
        })
        
        # Slide 2: Diagram
        slides.append({
            "title": "Process Flow",
            "type": "diagram",
            "content": {
                "nodes": [{"id": n.id, "label": n.label, "type": n.node_type, 
                          "x": n.x, "y": n.y} for n in self.nodes.values()],
                "edges": [{"source": e.source, "target": e.target, "label": e.label} 
                         for e in self.edges]
            }
        })
        
        # Slide 3: Statistics
        slides.append({
            "title": "Flow Statistics",
            "type": "bullet",
            "content": {
                "points": [
                    f"Total Nodes: {len(self.nodes)}",
                    f"Total Edges: {len(self.edges)}",
                    f"Start Points: {sum(1 for n in self.nodes.values() if n.node_type == 'startpoint')}",
                    f"End Points: {sum(1 for n in self.nodes.values() if n.node_type == 'endpoint')}",
                    f"Decision Points: {sum(1 for n in self.nodes.values() if n.node_type == 'decision')}"
                ]
            }
        })
        
        return {"slides": slides}
    
    def create_pptx_file(self) -> bytes:
        """
        Generate a real PowerPoint file (binary PPTX format)
        
        Returns:
            Binary PPTX file content
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
            from io import BytesIO
        except ImportError:
            raise ImportError("python-pptx library is required. Install it with: pip install python-pptx")
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Slide 1: Title Slide
        title_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = "Flow Diagram"
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 102, 204)
        
        # Add subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        p = subtitle_frame.paragraphs[0]
        p.text = f"Nodes: {len(self.nodes)} | Edges: {len(self.edges)}"
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(89, 89, 89)
        
        # Slide 2: Diagram Overview
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = "Process Flow Diagram"
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 102, 204)
        
        # Draw nodes and edges as shapes
        self._draw_diagram_on_slide(slide, start_y=1.2)
        
        # Slide 3: Statistics
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = "Flow Statistics"
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 102, 204)
        
        # Add statistics
        stats_y = 1.2
        stats = [
            f"Total Nodes: {len(self.nodes)}",
            f"Total Edges: {len(self.edges)}",
            f"Start Points: {sum(1 for n in self.nodes.values() if n.node_type == 'startpoint')}",
            f"End Points: {sum(1 for n in self.nodes.values() if n.node_type == 'endpoint')}",
            f"Decision Points: {sum(1 for n in self.nodes.values() if n.node_type == 'decision')}",
        ]
        
        for stat in stats:
            stats_box = slide.shapes.add_textbox(Inches(1), Inches(stats_y), Inches(8), Inches(0.6))
            stats_frame = stats_box.text_frame
            p = stats_frame.paragraphs[0]
            p.text = "• " + stat
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(50, 50, 50)
            stats_y += 0.8
        
        # Slide 4: Node Details
        if self.nodes:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
            title_frame = title_box.text_frame
            p = title_frame.paragraphs[0]
            p.text = "Node Details"
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 102, 204)
            
            # Add nodes table
            nodes_y = 1.2
            for node_id, node in list(self.nodes.items())[:8]:  # Show first 8 nodes
                node_text = f"• {node.label} (Type: {node.node_type})"
                node_box = slide.shapes.add_textbox(Inches(1), Inches(nodes_y), Inches(8), Inches(0.5))
                node_frame = node_box.text_frame
                p = node_frame.paragraphs[0]
                p.text = node_text
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(50, 50, 50)
                nodes_y += 0.6
        
        # Slide 5: Edge Details (show all connections and their labels)
        if self.edges:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
            title_frame = title_box.text_frame
            p = title_frame.paragraphs[0]
            p.text = "Flow Connections (Edges)"
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 102, 204)
            
            # Add edges table
            edges_y = 1.2
            for idx, edge in enumerate(self.edges[:10]):  # Show first 10 edges
                source_label = self.nodes[edge.source].label if edge.source in self.nodes else edge.source
                target_label = self.nodes[edge.target].label if edge.target in self.nodes else edge.target
                edge_text = f"• {source_label} → {target_label}"
                if edge.label:
                    edge_text += f" [{edge.label}]"
                
                edge_box = slide.shapes.add_textbox(Inches(1), Inches(edges_y), Inches(8), Inches(0.5))
                edge_frame = edge_box.text_frame
                edge_frame.word_wrap = True
                p = edge_frame.paragraphs[0]
                p.text = edge_text
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(50, 50, 50)
                edges_y += 0.6
        
        # Save to bytes
        output = BytesIO()
        prs.save(output)
        output.seek(0)
        return output.getvalue()
    
    def _draw_diagram_on_slide(self, slide, start_y: float = 1.2):
        """
        Draw the diagram nodes and edges on a PowerPoint slide
        
        Args:
            slide: PowerPoint slide object
            start_y: Starting Y position for drawing
        """
        from pptx.util import Inches, Pt
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.dml.color import RGBColor
        from pptx.enum.dml import MSO_LINE
        
        # Color map for node types
        colors = {
            "startpoint": RGBColor(0, 153, 0),      # Green
            "endpoint": RGBColor(204, 0, 0),        # Red
            "decision": RGBColor(255, 153, 0),      # Orange
            "process": RGBColor(0, 102, 204)        # Blue
        }
        
        if not self.nodes:
            return
        
        # Draw up to 15 nodes in a grid and track their positions for edge drawing
        nodes_list = list(self.nodes.values())[:15]
        cols = 4
        x_start = Inches(0.5)
        y_pos = Inches(start_y)
        
        # Dictionary to store node positions on slide for edge drawing
        node_positions = {}
        shape_objects = {}
        
        # First pass: Draw all nodes and record their positions
        for idx, node in enumerate(nodes_list):
            col = idx % cols
            row = idx // cols
            
            x = x_start + Inches(col * 2.3)
            y = y_pos + Inches(row * 1.3)
            
            # Determine shape based on node type
            if node.node_type == "startpoint" or node.node_type == "endpoint":
                shape_type = MSO_SHAPE.OVAL
                width, height = Inches(1.0), Inches(0.5)
            elif node.node_type == "decision":
                shape_type = MSO_SHAPE.DIAMOND
                width, height = Inches(1.0), Inches(0.7)
            else:
                shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
                width, height = Inches(1.2), Inches(0.5)
            
            # Add shape
            shape = slide.shapes.add_shape(
                shape_type,
                x, y, width, height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = colors.get(node.node_type, RGBColor(0, 102, 204))
            shape.line.color.rgb = RGBColor(50, 50, 50)
            shape.line.width = Pt(1.5)
            
            # Add label - preserve full descriptive text
            text_frame = shape.text_frame
            text_frame.word_wrap = True
            text_frame.margin_bottom = Inches(0.05)
            text_frame.margin_top = Inches(0.05)
            text_frame.margin_left = Inches(0.05)
            text_frame.margin_right = Inches(0.05)
            
            p = text_frame.paragraphs[0]
            p.text = node.label  # Use full label, not truncated
            p.font.size = Pt(9)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.bold = True
            p.alignment = 1  # Center alignment
            
            # Store position: center of shape for edge drawing
            # Convert to inches values (floats) for proper arithmetic
            center_x_inches = x.inches + width.inches / 2
            center_y_inches = y.inches + height.inches / 2
            node_positions[node.id] = (center_x_inches, center_y_inches, x, y, width, height, shape)
            shape_objects[node.id] = shape
        
        # Second pass: Draw edges between connected nodes
        for edge in self.edges:
            if edge.source not in node_positions or edge.target not in node_positions:
                continue
            
            source_pos = node_positions[edge.source]
            target_pos = node_positions[edge.target]
            
            # node_positions stores inch values as floats, not Inches objects
            source_center_x, source_center_y = source_pos[0], source_pos[1]
            target_center_x, target_center_y = target_pos[0], target_pos[1]
            
            # Draw connector line (arrow) - convert float inches back to Inches objects
            connector = slide.shapes.add_connector(1, Inches(source_center_x), Inches(source_center_y), 
                                                   Inches(target_center_x), Inches(target_center_y))
            connector.line.color.rgb = RGBColor(80, 80, 80)
            connector.line.width = Pt(1.5)
            
            # Add edge label if present
            if edge.label:
                # Calculate midpoint for label placement (source and target are already in inches as floats)
                mid_x = Inches((source_center_x + target_center_x) / 2)
                mid_y = Inches((source_center_y + target_center_y) / 2)
                
                # Add label text box
                label_box = slide.shapes.add_textbox(
                    mid_x - Inches(0.5), 
                    mid_y - Inches(0.15),
                    Inches(1.0), 
                    Inches(0.3)
                )
                label_frame = label_box.text_frame
                label_frame.word_wrap = True
                label_p = label_frame.paragraphs[0]
                label_p.text = edge.label
                label_p.font.size = Pt(8)
                label_p.font.color.rgb = RGBColor(80, 80, 80)
                label_p.font.italic = True
                label_p.alignment = 1  # Center
